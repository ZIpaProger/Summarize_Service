from typing import Optional


class Parser:
    """Класс-обёртка для извлечения текста из разных форматов файлов.

    Поддерживает: .txt, .docx, .pdf, .pptx
    """

    @staticmethod
    def from_path(path: str) -> str:
        ext = path.split('.')[-1].lower()

        if ext == "txt":
            with open(path, "r", encoding="utf-8") as f:
                return f.read()

        if ext == "docx":
            import docx
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == "pdf":
            import pdfplumber
            text = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text.append(page.extract_text())
            return "\n".join(filter(None, text))

        if ext == "pptx":
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)

        raise ValueError(f"Формат '{ext}' не поддерживается.")
